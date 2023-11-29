import os
import flask
from llama_index import GPTSimpleVectorIndex, SimpleDirectoryReader, MockLLMPredictor, ServiceContext, MockEmbedding
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def preprocess(directory):
    # iterate over files in
    # that directory
    for filename in os.listdir(directory):
        fullPath = os.path.join(directory, filename)
        # checking if it is a file
        if os.path.isfile(fullPath):
            extension = os.path.splitext(filename)[1]
            if extension == '.ppt':
                print('.ppts should be converted to .pptx file extensions')
                continue;
            if extension == '.pptx':
                prs = Presentation(os.path.join(fullPath))  # name of your pptx
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                            pic = shape._element
                            pic.getparent().remove(pic)
                prs.save(fullPath)  # save changes

def build_mock(directory):

    # Loop through files in directory, extract text and summarize

    preprocess(directory)

    documents = SimpleDirectoryReader(directory).load_data()
    llm_predictor = MockLLMPredictor(max_tokens=256)
    embed_model = MockEmbedding(embed_dim=1536)
    service_context = ServiceContext.from_defaults(llm_predictor=llm_predictor, embed_model=embed_model)
    index = GPTSimpleVectorIndex.from_documents(documents, service_context=service_context)

    response = flask.make_response('Mock Built and printed tokens')
    response.status_code = 200

    return response

def build(directory):

    # remove some ppt images that dont parse well.
    preprocess(directory)

    documents = SimpleDirectoryReader(directory).load_data()
    # Concatenate all the documents into a single text string

    # !!!!! API CALL HAPPENS HERE CAREFUL.
    index = GPTSimpleVectorIndex.from_documents(documents)

    if not os.path.exists(directory + "/results"):
        os.makedirs(directory + "/results")
        print("Successfully created " + directory)
    else:
        print("Directory already exists. Writing index.json")

    index.save_to_disk(directory + '/results/index.json')
    response = flask.make_response('built an index.json result.')
    response.status_code = 200

    return response